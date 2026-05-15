"""Generate DOCX (Word) and PPTX (slide deck) for Digit Memory v2.

Outputs:
    design/digit_memory_case_study.docx   Word document
    design/digit_memory_deck.pptx          Slide deck (~10 slides)
"""
from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Pt, RGBColor, Inches, Cm
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt, Emu
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

HERE = Path(__file__).parent
ARCH_IMG = HERE / "architecture.png"

# Palette
INK = RGBColor(0x13, 0x13, 0x12)
INK_SOFT = RGBColor(0x3A, 0x3A, 0x37)
INK_MUTE = RGBColor(0x8A, 0x8A, 0x85)
ACCENT = RGBColor(0xB8, 0x33, 0x1F)

P_INK = PRGBColor(0x13, 0x13, 0x12)
P_INK_SOFT = PRGBColor(0x3A, 0x3A, 0x37)
P_INK_MUTE = PRGBColor(0x8A, 0x8A, 0x85)
P_ACCENT = PRGBColor(0xB8, 0x33, 0x1F)
P_BG = PRGBColor(0xF2, 0xEF, 0xE6)


# ---------- DOCX ----------
def _set_run(run, font_name="Inter", size=11, color=INK, bold=False, italic=False):
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.bold = bold
    run.italic = italic
    # ensure east-asian font fallback set
    rPr = run._element.get_or_add_rPr()
    rFonts = rPr.find(qn("w:rFonts"))
    if rFonts is None:
        rFonts = OxmlElement("w:rFonts")
        rPr.append(rFonts)
    rFonts.set(qn("w:ascii"), font_name)
    rFonts.set(qn("w:hAnsi"), font_name)


def _para(doc, text, font="Helvetica", size=11, color=INK,
          bold=False, italic=False, align=None, space_after=6, space_before=0):
    p = doc.add_paragraph()
    p.paragraph_format.space_after = Pt(space_after)
    p.paragraph_format.space_before = Pt(space_before)
    if align is not None:
        p.alignment = align
    r = p.add_run(text)
    _set_run(r, font_name=font, size=size, color=color,
             bold=bold, italic=italic)
    return p


def _hr(doc):
    p = doc.add_paragraph()
    pPr = p._p.get_or_add_pPr()
    pBdr = OxmlElement("w:pBdr")
    bottom = OxmlElement("w:bottom")
    bottom.set(qn("w:val"), "single")
    bottom.set(qn("w:sz"), "4")
    bottom.set(qn("w:space"), "1")
    bottom.set(qn("w:color"), "131312")
    pBdr.append(bottom)
    pPr.append(pBdr)
    p.paragraph_format.space_after = Pt(6)


def build_docx():
    doc = Document()
    # page setup
    section = doc.sections[0]
    section.top_margin    = Inches(0.85)
    section.bottom_margin = Inches(0.85)
    section.left_margin   = Inches(0.90)
    section.right_margin  = Inches(0.90)

    # cover
    _para(doc, "CASE STUDY  /  N. JAIN  /  ML_PROJECTS",
          font="Menlo", size=9, color=INK_MUTE, space_after=24)
    _para(doc, "DIGIT MEMORY", font="Helvetica", size=48,
          color=INK, bold=True, space_after=0)
    _para(doc, "A study in retrieval as classification.",
          font="Georgia", size=14, color=INK_SOFT, italic=True,
          space_after=4)
    _para(doc, "Built three ways. No model trained.",
          font="Georgia", size=14, color=INK_SOFT, italic=True,
          space_after=20)
    _hr(doc)

    _para(doc, "98.4%", font="Helvetica", size=72,
          color=ACCENT, bold=True, space_after=0)
    _para(doc, "HELD-OUT ACCURACY",
          font="Menlo", size=10, color=INK, bold=True, space_after=2)
    _para(doc,
          "1-NN on a stratified 75/25 split of the sklearn digits dataset.",
          font="Menlo", size=9, color=INK_SOFT, space_after=18)
    _hr(doc)

    # Section 1 — Hypothesis
    _para(doc, "01  THE HYPOTHESIS",
          font="Menlo", size=10, color=ACCENT, bold=True,
          space_before=12, space_after=10)
    _para(doc, "Hypothesis", font="Helvetica", size=30,
          color=INK, bold=True, space_after=14)
    _para(doc,
          "The fastest model is the one you don't need. Most ML projects "
          "start with what model to use. This one started with a different "
          "question.",
          font="Georgia", size=14, color=INK, italic=True, space_after=12)
    _para(doc,
          "What if you just stored the data, organized it well, and looked "
          "things up? For a dataset of 1797 small samples in 64 dimensions, "
          "the answer turns out to be: that's enough.",
          font="Helvetica", size=11)
    _para(doc,
          "Two query patterns drive the design. An exact match — the input "
          "is byte-identical to something seen — wants a hash table. O(1) "
          "lookup. There is no smarter answer.",
          font="Helvetica", size=11)
    _para(doc,
          "An approximate match — the input is close but not identical — "
          "wants a spatial index. Hashing breaks: change one bit and the "
          "hash changes completely. We need a structure that knows about "
          "geometric closeness.",
          font="Helvetica", size=11)
    _para(doc,
          "The combined system tries the cheap lookup first and falls back "
          "to the expensive one only when needed. That is the entire idea.",
          font="Helvetica", size=11)

    # Section 2 — Architecture
    _para(doc, "02  ARCHITECTURE",
          font="Menlo", size=10, color=ACCENT, bold=True,
          space_before=14, space_after=10)
    _para(doc, "Two-tier lookup", font="Helvetica", size=30,
          color=INK, bold=True, space_after=10)

    if ARCH_IMG.exists():
        doc.add_picture(str(ARCH_IMG), width=Inches(6.4))
        # caption
        _para(doc, "FIG.01  HASH FIRST · TREE ON MISS",
              font="Menlo", size=9, color=INK_MUTE,
              align=WD_ALIGN_PARAGRAPH.CENTER, space_after=12)

    _para(doc,
          "Most production lookup systems quietly do this — caches in front "
          "of expensive indexes, content-addressed stores in front of "
          "full-text search. The same pattern applies cleanly to "
          "classification when the data fits in memory.",
          font="Helvetica", size=11)

    # Section 3 — Numbers
    _para(doc, "03  NUMBERS",
          font="Menlo", size=10, color=ACCENT, bold=True,
          space_before=14, space_after=10)
    _para(doc, "Numbers", font="Helvetica", size=30,
          color=INK, bold=True, space_after=8)

    metrics = [
        ("HELD-OUT ACCURACY", "98.4%",
         "stratified 75/25 split  ·  1-NN  ·  no training", True),
        ("EXACT LOOKUP", "~3 microseconds",
         "Python dict, bytes(features) -> label", False),
        ("NEAREST NEIGHBOR", "~50 microseconds",
         "BallTree query, k=1, 1797 samples", False),
    ]
    for top, big, sub, accent in metrics:
        col = ACCENT if accent else INK
        _para(doc, top, font="Menlo", size=10, color=col,
              bold=True, space_after=2)
        _para(doc, big, font="Helvetica", size=28, color=col,
              bold=True, space_after=2)
        _para(doc, sub, font="Menlo", size=9, color=INK_SOFT,
              space_after=10)

    # Section 4 — Robustness
    _para(doc, "04  ROBUSTNESS",
          font="Menlo", size=10, color=ACCENT, bold=True,
          space_before=14, space_after=10)
    _para(doc, "Robustness", font="Helvetica", size=30,
          color=INK, bold=True, space_after=8)
    _para(doc,
          "Accuracy on the held-out test set with injected Gaussian noise.",
          font="Georgia", size=12, color=INK_SOFT, italic=True,
          space_after=8)
    # table
    table = doc.add_table(rows=1, cols=2)
    table.style = "Light Grid Accent 1"
    hdr = table.rows[0].cells
    hdr[0].text = "Noise sigma (pixel units)"
    hdr[1].text = "Held-out accuracy"
    rows = [
        ("0.0",  "98.4%"),
        ("0.5",  "98.2%"),
        ("1.0",  "98.2%"),
        ("2.0",  "98.4%"),
        ("3.0",  "97.3%"),
        ("4.0",  "96.4%"),
        ("5.0",  "93.8%"),
        ("7.0",  "84.7%"),
        ("10.0", "67.1%"),
    ]
    for s, a in rows:
        row = table.add_row().cells
        row[0].text = s
        row[1].text = a
    _para(doc, "", space_after=4)
    _para(doc,
          "Smooth degradation, no cliff. That is what robust retrieval looks like.",
          font="Helvetica", size=11, italic=True, color=INK_SOFT)

    # Section 5 — Implementations
    _para(doc, "05  THREE IMPLEMENTATIONS",
          font="Menlo", size=10, color=ACCENT, bold=True,
          space_before=14, space_after=10)
    _para(doc, "Three ways", font="Helvetica", size=30,
          color=INK, bold=True, space_after=10)
    _para(doc,
          "Writing the same system three times is the closest I have found "
          "to actually understanding it.",
          font="Georgia", size=13, color=INK_SOFT, italic=True,
          space_after=10)

    impls = [
        ("PYTHON",  "digits_memory.py",
         "Production-style. SQLite persistence. sklearn KDTree.",
         "Use when you want the cleanest summary."),
        ("JUPYTER", "digits_study_guide.ipynb",
         "43 cells. Pandas DataFrames. Three NN backends benchmarked.",
         "Use when you want to learn the project hands-on."),
        ("C",       "digits_memory.c",
         "Zero dependencies. FNV-1a open-addressed hash. Brute-force NN.",
         "Use when you want to understand every byte."),
    ]
    for title, file, desc, when in impls:
        _para(doc, title, font="Helvetica", size=20,
              color=INK, bold=True, space_after=2)
        _para(doc, file, font="Menlo", size=10, color=ACCENT,
              bold=True, space_after=4)
        _para(doc, desc, font="Helvetica", size=11, color=INK,
              space_after=2)
        _para(doc, when, font="Georgia", size=11, color=INK_SOFT,
              italic=True, space_after=12)

    # Section 6 — Takeaways
    _para(doc, "06  TAKEAWAYS",
          font="Menlo", size=10, color=ACCENT, bold=True,
          space_before=14, space_after=10)
    _para(doc, "What I learned", font="Helvetica", size=30,
          color=INK, bold=True, space_after=10)

    takeaways = [
        ("Retrieval is a criminally underused baseline.",
         "Before tuning a model, ask whether you can just remember the "
         "answer. On small-to-medium data the answer is often yes."),
        ("Match the structure to the query distribution.",
         "Exact lookups want a hash. Approximate lookups want a tree. "
         "Production systems should usually do both."),
        ("Tree NN is not always faster than brute force.",
         "At 1797 samples and 64 dimensions, vectorized numpy beat KDTree "
         "and BallTree. Asymptotics flip at larger N. Always benchmark."),
        ("Hash key choice matters.",
         "`array.tobytes()` is the right key for numpy arrays. "
         "`tuple(array)` is slower. The difference is measurable."),
        ("Writing the same system three ways compounds understanding.",
         "Python lets you build it. C makes you understand it."),
    ]
    for i, (title, body) in enumerate(takeaways, 1):
        _para(doc, f"{i:02d}.  {title}",
              font="Helvetica", size=12, color=INK, bold=True,
              space_after=2)
        _para(doc, body, font="Helvetica", size=11, color=INK_SOFT,
              space_after=10)

    # Colophon
    _para(doc, "07  COLOPHON",
          font="Menlo", size=10, color=ACCENT, bold=True,
          space_before=14, space_after=10)
    _para(doc, "Colophon", font="Helvetica", size=30,
          color=INK, bold=True, space_after=10)
    coloph = [
        ("Design philosophy",  "Computational Quiet"),
        ("Display typeface",   "Big Shoulders Bold (Helvetica in this DOCX)"),
        ("Body typeface",      "Instrument Sans (Helvetica in this DOCX)"),
        ("Mono typeface",      "Geist Mono (Menlo in this DOCX)"),
        ("Repository",         "github.com/nimb-ou/ML_projects"),
        ("License",            "MIT"),
        ("Author",             "Nimit Jain"),
    ]
    for k, v in coloph:
        _para(doc, f"{k:<22}  {v}",
              font="Menlo", size=10, color=INK, space_after=2)

    doc.save(HERE / "digit_memory_case_study.docx")


# ---------- PPTX ----------
def _add_text(slide, x, y, w, h, text, font_name="Helvetica",
              size=18, color=P_INK, bold=False, italic=False,
              align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font_name
    r.font.size = PPt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    return tb


def _add_rule(slide, x, y, w, color=P_INK, weight=0.5):
    line = slide.shapes.add_connector(1, x, y, x + w, y)
    line.line.color.rgb = color
    line.line.width = Emu(int(weight * 9525))


def _fill_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = P_BG


def build_pptx():
    prs = Presentation()
    # 16:9 widescreen 13.333 x 7.5 inches
    prs.slide_width  = PInches(13.333)
    prs.slide_height = PInches(7.5)
    blank = prs.slide_layouts[6]

    SW = prs.slide_width
    SH = prs.slide_height
    M = PInches(0.55)

    # --- Slide 1: cover ---
    s = prs.slides.add_slide(blank)
    _fill_bg(s)
    _add_text(s, M, PInches(0.45), PInches(8), PInches(0.3),
              "DIGIT MEMORY  /  v2  /  RETRIEVAL AS A BASELINE",
              font_name="Menlo", size=10, color=P_INK_MUTE)
    _add_text(s, M, PInches(1.0), PInches(11), PInches(2.0),
              "Digit Memory", font_name="Helvetica", size=88,
              bold=True, color=P_INK)
    _add_text(s, M, PInches(2.4), PInches(11), PInches(0.6),
              "A study in retrieval as classification.",
              font_name="Georgia", size=22, italic=True, color=P_INK_SOFT)
    _add_rule(s, M, PInches(3.3), SW - 2 * M)
    _add_text(s, M, PInches(3.6), PInches(8), PInches(2.0),
              "98.4%", font_name="Helvetica", size=160, bold=True,
              color=P_ACCENT)
    _add_text(s, M, PInches(5.7), PInches(11), PInches(0.4),
              "HELD-OUT ACCURACY",
              font_name="Menlo", size=12, bold=True, color=P_INK)
    _add_text(s, M, PInches(6.0), PInches(11), PInches(0.35),
              "1-NN on a stratified 75/25 split of the sklearn digits dataset.",
              font_name="Menlo", size=10, color=P_INK_SOFT)
    _add_text(s, M, PInches(6.9), PInches(11), PInches(0.3),
              "Nimit Jain  /  github.com/nimb-ou/ML_projects",
              font_name="Menlo", size=9, color=P_INK_MUTE)

    # --- Slide 2: the question ---
    s = prs.slides.add_slide(blank)
    _fill_bg(s)
    _add_text(s, M, PInches(0.45), PInches(8), PInches(0.3),
              "§ 01  ·  THE QUESTION",
              font_name="Menlo", size=10, color=P_INK_MUTE)
    _add_rule(s, M, PInches(0.85), SW - 2 * M)
    _add_text(s, M, PInches(1.4), SW - 2 * M, PInches(5),
              "The fastest model is the one\nyou don't need.",
              font_name="Helvetica", size=72, bold=True, color=P_INK)
    _add_text(s, M, PInches(4.0), SW - 2 * M, PInches(3),
              "Most ML projects start with what model to use.\n"
              "This one started with a different question:\n"
              "what if you just stored the data, organized it well,\n"
              "and looked things up?",
              font_name="Georgia", size=26, italic=True, color=P_INK_SOFT)

    # --- Slide 3: architecture ---
    s = prs.slides.add_slide(blank)
    _fill_bg(s)
    _add_text(s, M, PInches(0.45), PInches(8), PInches(0.3),
              "§ 02  ·  ARCHITECTURE",
              font_name="Menlo", size=10, color=P_INK_MUTE)
    _add_rule(s, M, PInches(0.85), SW - 2 * M)
    _add_text(s, M, PInches(1.0), PInches(11), PInches(1.0),
              "Two-tier lookup",
              font_name="Helvetica", size=44, bold=True, color=P_INK)
    if ARCH_IMG.exists():
        s.shapes.add_picture(str(ARCH_IMG),
                             PInches(1.0), PInches(2.1),
                             width=PInches(11.3))
    _add_text(s, M, PInches(6.95), PInches(11), PInches(0.3),
              "Hash first. Tree on miss. Return whichever answer arrives.",
              font_name="Georgia", size=14, italic=True, color=P_INK_SOFT)

    # --- Slide 4: numbers ---
    s = prs.slides.add_slide(blank)
    _fill_bg(s)
    _add_text(s, M, PInches(0.45), PInches(8), PInches(0.3),
              "§ 03  ·  NUMBERS",
              font_name="Menlo", size=10, color=P_INK_MUTE)
    _add_rule(s, M, PInches(0.85), SW - 2 * M)
    col_w = (SW - 2 * M) / 3
    headers = [
        ("HELD-OUT ACCURACY", "98.4%", "stratified 75/25 split", True),
        ("EXACT LOOKUP",      "3 us",  "Python dict, O(1)",       False),
        ("NEAREST NEIGHBOR",  "50 us", "BallTree, k=1",           False),
    ]
    for i, (top, big, sub, accent) in enumerate(headers):
        col = P_ACCENT if accent else P_INK
        cx = M + i * col_w
        _add_text(s, cx, PInches(1.3), col_w, PInches(0.4),
                  top, font_name="Menlo", size=11, bold=True, color=col)
        _add_text(s, cx, PInches(1.9), col_w, PInches(3.0),
                  big, font_name="Helvetica", size=120, bold=True, color=col)
        _add_text(s, cx, PInches(5.4), col_w, PInches(0.4),
                  sub, font_name="Menlo", size=11, color=P_INK_SOFT)
    _add_rule(s, M, PInches(6.4), SW - 2 * M)
    _add_text(s, M, PInches(6.55), SW - 2 * M, PInches(0.5),
              "No training. No gradients. No epochs.",
              font_name="Georgia", size=18, italic=True, color=P_INK_SOFT)

    # --- Slide 5: implementations ---
    s = prs.slides.add_slide(blank)
    _fill_bg(s)
    _add_text(s, M, PInches(0.45), PInches(8), PInches(0.3),
              "§ 04  ·  IMPLEMENTATIONS",
              font_name="Menlo", size=10, color=P_INK_MUTE)
    _add_rule(s, M, PInches(0.85), SW - 2 * M)
    _add_text(s, M, PInches(1.0), PInches(11), PInches(1.0),
              "Three ways",
              font_name="Helvetica", size=44, bold=True, color=P_INK)
    _add_text(s, M, PInches(2.1), SW - 2 * M, PInches(0.5),
              "Writing the same system three times is the closest I have found to actually understanding it.",
              font_name="Georgia", size=18, italic=True, color=P_INK_SOFT)
    rows = [
        ("01", "PYTHON",  "digits_memory.py",
         "Production-style. SQLite persistence. sklearn KDTree."),
        ("02", "JUPYTER", "digits_study_guide.ipynb",
         "43 cells. Pandas DataFrames. Three NN backends benchmarked."),
        ("03", "C",       "digits_memory.c",
         "Zero dependencies. FNV-1a open-addressed hash. Brute-force NN."),
    ]
    y0 = PInches(3.2)
    for i, (num, name, file, desc) in enumerate(rows):
        y = y0 + PInches(1.2 * i)
        _add_text(s, M, y, PInches(0.6), PInches(0.5), num,
                  font_name="Menlo", size=12, color=P_INK_MUTE)
        _add_text(s, M + PInches(0.7), y - PInches(0.05),
                  PInches(3.5), PInches(0.7), name,
                  font_name="Helvetica", size=32, bold=True, color=P_INK)
        _add_text(s, SW - M - PInches(4.5), y + PInches(0.07),
                  PInches(4.5), PInches(0.5), file,
                  font_name="Menlo", size=14, bold=True, color=P_ACCENT,
                  align=PP_ALIGN.RIGHT)
        _add_text(s, M + PInches(0.7), y + PInches(0.55),
                  PInches(11), PInches(0.4), desc,
                  font_name="Helvetica", size=14, color=P_INK_SOFT)
        _add_rule(s, M, y + PInches(1.05), SW - 2 * M, weight=0.3)

    # --- Slide 6: robustness ---
    s = prs.slides.add_slide(blank)
    _fill_bg(s)
    _add_text(s, M, PInches(0.45), PInches(8), PInches(0.3),
              "§ 05  ·  ROBUSTNESS",
              font_name="Menlo", size=10, color=P_INK_MUTE)
    _add_rule(s, M, PInches(0.85), SW - 2 * M)
    _add_text(s, M, PInches(1.0), PInches(11), PInches(1.0),
              "How robust?",
              font_name="Helvetica", size=44, bold=True, color=P_INK)
    _add_text(s, M, PInches(2.1), SW - 2 * M, PInches(0.6),
              "Accuracy on the held-out test set with injected Gaussian noise.",
              font_name="Georgia", size=18, italic=True, color=P_INK_SOFT)

    # three-column numbers
    bullets = [
        ("sigma = 0",  "98.4%"),
        ("sigma = 5",  "93.8%"),
        ("sigma = 10", "67.1%"),
    ]
    col_w = (SW - 2 * M) / 3
    for i, (k, v) in enumerate(bullets):
        cx = M + i * col_w
        _add_text(s, cx, PInches(3.3), col_w, PInches(0.4),
                  k, font_name="Menlo", size=14, bold=True, color=P_INK_MUTE)
        _add_text(s, cx, PInches(3.7), col_w, PInches(2.0),
                  v, font_name="Helvetica", size=88, bold=True, color=P_INK)
    _add_rule(s, M, PInches(6.4), SW - 2 * M)
    _add_text(s, M, PInches(6.55), SW - 2 * M, PInches(0.6),
              "Smooth degradation. No cliff. No catastrophic forgetting.",
              font_name="Georgia", size=18, italic=True, color=P_INK_SOFT)

    # --- Slide 7: takeaways ---
    s = prs.slides.add_slide(blank)
    _fill_bg(s)
    _add_text(s, M, PInches(0.45), PInches(8), PInches(0.3),
              "§ 06  ·  TAKEAWAYS",
              font_name="Menlo", size=10, color=P_INK_MUTE)
    _add_rule(s, M, PInches(0.85), SW - 2 * M)
    _add_text(s, M, PInches(1.0), PInches(11), PInches(1.0),
              "What I learned",
              font_name="Helvetica", size=44, bold=True, color=P_INK)

    takes = [
        ("Retrieval is an underused baseline.",
         "On small-to-medium data, just remembering the answer often works."),
        ("Match the structure to the query distribution.",
         "Exact lookups -> hash. Approximate -> tree. Production wants both."),
        ("Always benchmark before you assume.",
         "At 1797 x 64, vectorized brute force beat KDTree and BallTree."),
        ("Writing it three times compounds understanding.",
         "Python lets you build it. C makes you understand it."),
    ]
    y0 = PInches(2.1)
    for i, (head, body) in enumerate(takes):
        y = y0 + PInches(1.15 * i)
        _add_text(s, M, y, PInches(0.6), PInches(0.5), f"{i+1:02d}",
                  font_name="Menlo", size=12, color=P_ACCENT, bold=True)
        _add_text(s, M + PInches(0.7), y - PInches(0.05),
                  SW - 2 * M - PInches(0.7), PInches(0.5),
                  head, font_name="Helvetica", size=20, bold=True,
                  color=P_INK)
        _add_text(s, M + PInches(0.7), y + PInches(0.4),
                  SW - 2 * M - PInches(0.7), PInches(0.5),
                  body, font_name="Helvetica", size=14, color=P_INK_SOFT)
        _add_rule(s, M, y + PInches(1.0), SW - 2 * M, weight=0.3)

    # --- Slide 8: closing ---
    s = prs.slides.add_slide(blank)
    _fill_bg(s)
    _add_text(s, M, PInches(0.45), PInches(8), PInches(0.3),
              "§ 07  ·  END",
              font_name="Menlo", size=10, color=P_INK_MUTE)
    _add_rule(s, M, PInches(0.85), SW - 2 * M)
    _add_text(s, M, PInches(1.6), SW - 2 * M, PInches(2),
              "Try it.\nThen try beating it.",
              font_name="Helvetica", size=80, bold=True, color=P_INK)
    _add_rule(s, M, PInches(5.4), SW - 2 * M)
    _add_text(s, M, PInches(5.6), SW - 2 * M, PInches(0.5),
              "github.com/nimb-ou/ML_projects",
              font_name="Menlo", size=22, bold=True, color=P_ACCENT)
    _add_text(s, M, PInches(6.2), SW - 2 * M, PInches(0.5),
              "MIT  ·  Code, blog, notebooks, slides — all in the repo.",
              font_name="Menlo", size=12, color=P_INK_SOFT)
    _add_text(s, M, PInches(6.9), SW - 2 * M, PInches(0.3),
              "Nimit Jain  ·  2026",
              font_name="Georgia", size=14, italic=True, color=P_INK_MUTE)

    prs.save(HERE / "digit_memory_deck.pptx")


if __name__ == "__main__":
    print("[1/2] building DOCX ...")
    build_docx()
    print("[2/2] building PPTX ...")
    build_pptx()
    print("Done.")
